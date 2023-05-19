# you will also need to install Pillow: https://pypi.org/project/Pillow:
# pip install Pillow

# added to patch bug in 3.11.3 where the following error is thrown:
# Traceback (most recent call last):
#  File "C:\Users\njdem\AppData\Local\Packages\PythonSoftwareFoundation.Python.3.11_qbz5n2kfra8p0\LocalCache\local-packages\Python311\site-packages\pptx\compat\__init__.py", line 10, in <module>
#    Container = collections.abc.Container
#                ^^^^^^^^^^^^^^^
# AttributeError: module 'collections' has no attribute 'abc'

import collections
import collections.abc

from pptx import Presentation #install via pip install python-pptx

import os # Use PowerPoint to open our PPT


# Creating presentation object
ppt = Presentation()

# slide types
# 0 Title and subtitle
# 1 Title and content
# 2 Section header
# 3 Two content
# 4 Comparison
# 5 Title only 
# 6 Blank
# 7 Content with caption
# 8 Image with caption

# Creating first slide layout
first_slide_layout = ppt.slide_layouts[0] # 0 is the slide type

# Add new slide to the deck
slide = ppt.slides.add_slide(first_slide_layout)

# Add a slide title
slide.shapes.title.text = "Created by Phthon using python-pptx"

second_slide_layout = ppt.slide_layouts[1] # 1 is the slide type

slide = ppt.slides.add_slide(second_slide_layout)

slide.shapes.title.text = "Second Slide"

# Save the new PowerPoint deck under folder where CreateDeck.py is ran
ppt.save("PythonPPT.pptx")

print("New PPT created, opening....")

os.startfile("PythonPPT.pptx")
